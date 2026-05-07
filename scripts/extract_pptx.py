import collections 
import collections.abc
from pptx import Presentation

prs = Presentation('C:\\Data\\PlaySmart\\Week8_chapter12_ID6e.pptx')
text_runs = []

for slide in prs.slides:
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            text_runs.append(shape.text)

with open('C:\\Data\\PlaySmart\\slides_text.txt', 'w', encoding='utf-8') as f:
    f.write("\n".join(text_runs))
